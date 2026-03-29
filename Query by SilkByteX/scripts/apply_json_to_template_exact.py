from __future__ import annotations

import json
from pathlib import Path
from pptx import Presentation

TEMPLATE = Path('/Users/mobashirsifat/Downloads/BIG DATA PROJECT UPDATED 1.pptx')
CONTENT = Path('yelp_text_to_sql/data_journey_content.json')
OUT = Path('output/SILKBYTE_X_QUERY_Data_Journey_TEMPLATE_STYLE.pptx')


def text_shapes(slide):
    shapes = [sh for sh in slide.shapes if getattr(sh, 'has_text_frame', False)]
    return sorted(shapes, key=lambda sh: (sh.top, sh.left))


def set_text(shape, text: str):
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = text


def choose_targets(slide):
    ts = text_shapes(slide)
    non_empty = [s for s in ts if (s.text or '').strip()]
    # Keep tiny labels (logo/team names) untouched by preferring medium/large text boxes.
    candidates = [s for s in non_empty if s.width > 2500000 and s.height > 300000]
    if not candidates:
        candidates = non_empty
    title = None
    body = None
    chips = None
    # Highest candidate for title.
    if candidates:
        title = sorted(candidates, key=lambda s: (s.top, -s.width))[0]
    # Longest text candidate for body.
    if candidates:
        body = sorted(candidates, key=lambda s: len((s.text or '').strip()), reverse=True)[0]
    # Another candidate near top for chips.
    if len(candidates) > 1:
        near_top = sorted(candidates, key=lambda s: (s.top, s.left))
        for c in near_top:
            if c is not title and c is not body:
                chips = c
                break
    return title, chips, body


def apply_slide_content(slide, title_text: str, chips_text: str, body_text: str):
    title, chips, body = choose_targets(slide)
    if title is not None:
        set_text(title, title_text)
    if chips is not None and chips_text:
        set_text(chips, chips_text)
    if body is not None:
        set_text(body, body_text)


def main():
    content = json.loads(CONTENT.read_text(encoding='utf-8'))
    prs = Presentation(str(TEMPLATE))

    # Slide 1 (cover)
    s1 = prs.slides[0]
    ts1 = [s for s in text_shapes(s1) if (s.text or '').strip()]
    if ts1:
        ts1_sorted = sorted(ts1, key=lambda s: (s.top, s.left))
        set_text(ts1_sorted[0], content.get('page_title', ''))
        if len(ts1_sorted) > 1:
            set_text(ts1_sorted[1], content.get('page_subtitle', ''))

    # Slide 2 opening
    op = content.get('opening', {})
    if len(prs.slides) >= 2:
        apply_slide_content(
            prs.slides[1],
            op.get('title', ''),
            op.get('kicker', ''),
            f"{op.get('line_1','')}\n{op.get('line_2','')}\n{op.get('connect','')}",
        )

    # Slides 3..27 chapters
    chapters = content.get('chapters', [])
    for i, ch in enumerate(chapters):
        slide_idx = 2 + i
        if slide_idx >= len(prs.slides):
            break
        chips = '  •  '.join(ch.get('chips', []))
        body = f"{ch.get('line_1','')}\n{ch.get('line_2','')}\n{ch.get('connect','')}"
        apply_slide_content(prs.slides[slide_idx], ch.get('title', ''), chips, body)

    # Slide 28 closing
    closing_idx = 2 + len(chapters)
    cl = content.get('closing', {})
    if closing_idx < len(prs.slides):
        apply_slide_content(
            prs.slides[closing_idx],
            cl.get('title', ''),
            cl.get('kicker', ''),
            f"{cl.get('line_1','')}\n{cl.get('line_2','')}\n{cl.get('line_3','')}\n{cl.get('connect','')}",
        )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f'Created: {OUT.resolve()}')


if __name__ == '__main__':
    main()
