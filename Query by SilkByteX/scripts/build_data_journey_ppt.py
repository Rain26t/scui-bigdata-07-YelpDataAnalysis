from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(8, 11, 22)
CARD = RGBColor(17, 20, 36)
CARD_ALT = RGBColor(12, 16, 31)
TITLE = RGBColor(247, 240, 255)
BODY = RGBColor(239, 233, 250)
SUBTLE = RGBColor(190, 180, 211)
PINK = RGBColor(233, 132, 216)
CYAN = RGBColor(142, 226, 255)
CHIP_FILL = RGBColor(68, 42, 88)
WHITE = RGBColor(255, 255, 255)


def _safe_image_path(base: Path, raw: str | None) -> Path | None:
    if not raw:
        return None
    candidate = (base / raw).resolve()
    return candidate if candidate.exists() else None


def _set_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_blob(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def _add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = CARD, radius: int = 22):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(80, 73, 110)
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = radius / 100.0
    return shape


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = "Aptos",
    margin_left: float = 0.02,
    margin_right: float = 0.02,
    margin_top: float = 0.02,
    margin_bottom: float = 0.02,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    box.text_frame.margin_left = Inches(margin_left)
    box.text_frame.margin_right = Inches(margin_right)
    box.text_frame.margin_top = Inches(margin_top)
    box.text_frame.margin_bottom = Inches(margin_bottom)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color


def _add_chip(slide, text: str, left: float, top: float, width: float) -> None:
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.34),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = CHIP_FILL
    chip.line.color.rgb = RGBColor(110, 94, 145)
    chip.line.width = Pt(0.8)
    chip.adjustments[0] = 0.35
    _add_textbox(
        slide,
        text,
        left + 0.06,
        top + 0.02,
        width - 0.12,
        0.24,
        font_size=10,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Aptos",
    )


def _add_chip_row(slide, chips: list[str], left: float, top: float, max_width: float) -> float:
    x = left
    y = top
    row_height = 0.42
    for chip in chips:
        chip_width = min(max(0.7 + len(chip) * 0.07, 1.05), 2.2)
        if x + chip_width > left + max_width:
            x = left
            y += row_height
        _add_chip(slide, chip, x, y, chip_width)
        x += chip_width + 0.12
    return y + row_height


def _image_size(path: Path) -> tuple[int, int]:
    try:
        with Image.open(path) as img:
            return img.size
    except Exception:
        return (1600, 900)


def _add_picture_contain(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    img_w, img_h = _image_size(path)
    box_ratio = width / height
    img_ratio = img_w / img_h if img_h else 1
    if img_ratio > box_ratio:
        draw_w = width
        draw_h = width / img_ratio
        draw_left = left
        draw_top = top + (height - draw_h) / 2
    else:
        draw_h = height
        draw_w = height * img_ratio
        draw_top = top
        draw_left = left + (width - draw_w) / 2
    slide.shapes.add_picture(str(path), Inches(draw_left), Inches(draw_top), width=Inches(draw_w), height=Inches(draw_h))


def _add_media_card(slide, base: Path, slot: dict[str, Any] | None, left: float, top: float, width: float, height: float) -> None:
    _add_card(slide, left, top, width, height, fill=CARD_ALT)
    if not slot:
        return
    path = _safe_image_path(base, str(slot.get("url", "")).strip())
    if path:
        _add_picture_contain(slide, path, left + 0.16, top + 0.16, width - 0.32, height - 0.74)
    else:
        label = str(slot.get("label", "Visual")).strip() or "Visual"
        _add_textbox(
            slide,
            label,
            left + 0.3,
            top + 0.55,
            width - 0.6,
            0.8,
            font_size=16,
            color=BODY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    caption = str(slot.get("caption", "")).strip()
    if caption:
        _add_textbox(
            slide,
            caption,
            left + 0.24,
            top + height - 0.42,
            width - 0.48,
            0.28,
            font_size=9,
            color=SUBTLE,
            align=PP_ALIGN.CENTER,
        )


def _decorate(slide, base: Path) -> None:
    pink = base / "assets/data_journey/ppt/design_blob_magenta.png"
    blue = base / "assets/data_journey/ppt/design_blob_blue.png"
    _add_blob(slide, pink, -0.55, -0.15, 3.2, 3.2)
    _add_blob(slide, blue, 10.55, 4.75, 3.2, 3.2)


def _add_title_slide(prs: Presentation, base: Path, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _decorate(slide, base)
    _add_card(slide, 0.62, 0.62, 12.05, 6.2)
    _add_textbox(slide, "DATA JOURNEY", 1.0, 1.05, 2.6, 0.3, font_size=12, color=PINK, bold=True)
    _add_textbox(slide, title, 0.98, 1.55, 9.7, 1.25, font_size=28, color=TITLE, bold=True)
    _add_textbox(slide, subtitle, 1.0, 2.78, 9.8, 0.8, font_size=15, color=BODY)
    _add_textbox(
        slide,
        "Raw Yelp JSON -> HDFS -> Hive -> PySpark -> Zeppelin -> Query by SilkByteX",
        1.0,
        4.84,
        7.0,
        0.42,
        font_size=12,
        color=CYAN,
        bold=True,
    )
    hero = _safe_image_path(base, "assets/data_journey/generated/closing_value_chain_v2.png")
    if hero:
        _add_picture_contain(slide, hero, 8.5, 3.2, 3.25, 2.35)


def _add_chapter_slide(
    prs: Presentation,
    base: Path,
    kicker: str,
    title: str,
    chips: list[str],
    line_1: str,
    line_2: str,
    connect: str,
    slot_1: dict[str, Any] | None,
    slot_2: dict[str, Any] | None,
    *,
    index_label: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _decorate(slide, base)
    _add_card(slide, 0.48, 0.38, 12.35, 2.56)
    kicker_text = kicker if not index_label else f"{kicker}  |  {index_label}"
    _add_textbox(slide, kicker_text.upper(), 0.82, 0.66, 4.0, 0.25, font_size=11, color=PINK, bold=True)
    _add_textbox(slide, title, 0.8, 1.0, 11.6, 0.62, font_size=24, color=TITLE, bold=True)
    chip_end = _add_chip_row(slide, chips, 0.82, 1.63, 8.7)
    line_top = max(2.0, chip_end + 0.02)
    _add_textbox(slide, line_1, 0.82, line_top, 11.2, 0.32, font_size=12, color=BODY)
    _add_textbox(slide, line_2, 0.82, line_top + 0.34, 11.2, 0.32, font_size=12, color=BODY)
    _add_textbox(slide, connect, 0.82, line_top + 0.72, 11.2, 0.25, font_size=10, color=CYAN, bold=True)
    _add_media_card(slide, base, slot_1, 0.48, 3.18, 6.07, 3.72)
    _add_media_card(slide, base, slot_2, 6.78, 3.18, 6.07, 3.72)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    base = root / "yelp_text_to_sql"
    content_path = base / "data_journey_content.json"
    out_path = root / "output" / "SILKBYTE_X_QUERY_Website_Style.pptx"

    content = json.loads(content_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = str(content.get("page_title", "SILKBYTE X QUERY - Data Journey"))
    prs.core_properties.subject = "Website-style data journey presentation"
    prs.core_properties.author = "Codex"

    _add_title_slide(prs, base, str(content.get("page_title", "")), str(content.get("page_subtitle", "")))

    opening = dict(content.get("opening", {}))
    hero_slots = list(content.get("hero_slots", []))
    while len(hero_slots) < 2:
        hero_slots.append({})
    _add_chapter_slide(
        prs,
        base,
        str(opening.get("kicker", "")),
        str(opening.get("title", "")),
        ["Project Vision", "System Narrative", "Data Journey"],
        str(opening.get("line_1", "")),
        str(opening.get("line_2", "")),
        str(opening.get("connect", "")),
        hero_slots[0],
        hero_slots[1],
    )

    for index, chapter in enumerate(content.get("chapters", []), start=1):
        _add_chapter_slide(
            prs,
            base,
            str(chapter.get("kicker", "")),
            str(chapter.get("title", "")),
            [str(chip) for chip in chapter.get("chips", [])],
            str(chapter.get("line_1", "")),
            str(chapter.get("line_2", "")),
            str(chapter.get("connect", "")),
            dict(chapter.get("slot_1", {})),
            dict(chapter.get("slot_2", {})),
            index_label=str(index),
        )

    closing = dict(content.get("closing", {}))
    finale = dict(content.get("finale_slot", {}))
    _add_chapter_slide(
        prs,
        base,
        str(closing.get("kicker", "")),
        str(closing.get("title", "")),
        ["End-to-End", "Explainable Analytics", "Conversational Intelligence"],
        str(closing.get("line_1", "")),
        str(closing.get("line_2", "")),
        f"{closing.get('line_3', '')}  {closing.get('connect', '')}".strip(),
        finale,
        {
            "url": "assets/data_journey/generated/closing_value_chain_v2.png",
            "caption": str(finale.get("caption", "")),
        },
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Created: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
