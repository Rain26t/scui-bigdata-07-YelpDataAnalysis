from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(8, 11, 22)
CARD = RGBColor(17, 20, 36)
CARD_ALT = RGBColor(12, 16, 31)
TITLE = RGBColor(247, 240, 255)
BODY = RGBColor(239, 233, 250)
SUBTLE = RGBColor(190, 180, 211)
PINK = RGBColor(233, 132, 216)
CYAN = RGBColor(142, 226, 255)
CHIP_FILL = RGBColor(68, 42, 88)
WHITE = RGBColor(255, 255, 255)


def safe_image_path(base: Path, raw: str | None) -> Path | None:
    if not raw:
        return None
    candidate = (base / raw).resolve()
    return candidate if candidate.exists() else None


def set_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_blob(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = CARD):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(80, 73, 110)
    shape.line.width = Pt(1.0)
    shape.adjustments[0] = 0.12
    return shape


def add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    box.text_frame.margin_left = Inches(0.02)
    box.text_frame.margin_right = Inches(0.02)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Aptos"
    run.font.color.rgb = color


def add_bullet_lines(slide, items: list[str], left: float, top: float, width: float, height: float, *, size: int = 14) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.02)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = box.text_frame.paragraphs[0] if idx == 0 else box.text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.name = "Aptos"
        run.font.color.rgb = BODY


def add_chip(slide, text: str, left: float, top: float, width: float) -> None:
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.34),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = CHIP_FILL
    chip.line.color.rgb = RGBColor(110, 94, 145)
    chip.line.width = Pt(0.8)
    chip.adjustments[0] = 0.35
    add_textbox(slide, text, left + 0.05, top + 0.02, width - 0.1, 0.22, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_chip_row(slide, chips: list[str], left: float, top: float, max_width: float) -> float:
    x = left
    y = top
    for chip in chips:
        width = min(max(0.75 + len(chip) * 0.065, 1.2), 2.5)
        if x + width > left + max_width:
            x = left
            y += 0.42
        add_chip(slide, chip, x, y, width)
        x += width + 0.12
    return y + 0.42


def image_size(path: Path) -> tuple[int, int]:
    try:
        with Image.open(path) as img:
            return img.size
    except Exception:
        return (1600, 900)


def add_picture_contain(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    img_w, img_h = image_size(path)
    box_ratio = width / height
    img_ratio = img_w / img_h if img_h else 1
    if img_ratio > box_ratio:
        draw_w = width
        draw_h = width / img_ratio
        draw_left = left
        draw_top = top + (height - draw_h) / 2
    else:
        draw_h = height
        draw_w = height * img_ratio
        draw_top = top
        draw_left = left + (width - draw_w) / 2
    slide.shapes.add_picture(str(path), Inches(draw_left), Inches(draw_top), width=Inches(draw_w), height=Inches(draw_h))


def decorate(slide, base: Path) -> None:
    add_blob(slide, base / "assets/data_journey/ppt/design_blob_magenta.png", -0.55, -0.15, 3.2, 3.2)
    add_blob(slide, base / "assets/data_journey/ppt/design_blob_blue.png", 10.55, 4.75, 3.2, 3.2)


def title_block(slide, kicker: str, title: str, subtitle: str, chips: list[str] | None = None) -> float:
    add_textbox(slide, kicker.upper(), 0.82, 0.62, 4.0, 0.24, size=11, color=PINK, bold=True)
    add_textbox(slide, title, 0.8, 0.95, 11.4, 0.62, size=24, color=TITLE, bold=True)
    add_textbox(slide, subtitle, 0.82, 1.6, 11.5, 0.55, size=12, color=BODY)
    chip_bottom = 2.08
    if chips:
        chip_bottom = add_chip_row(slide, chips, 0.82, 2.12, 10.7)
    return chip_bottom


def add_visual_card(slide, base: Path, image_rel: str | None, caption: str, left: float, top: float, width: float, height: float) -> None:
    add_card(slide, left, top, width, height, fill=CARD_ALT)
    path = safe_image_path(base, image_rel)
    if path:
        add_picture_contain(slide, path, left + 0.14, top + 0.14, width - 0.28, height - 0.58)
    else:
        add_textbox(slide, "Visual", left + 0.2, top + 0.65, width - 0.4, 0.5, size=16, color=BODY, bold=True, align=PP_ALIGN.CENTER)
    if caption:
        add_textbox(slide, caption, left + 0.18, top + height - 0.28, width - 0.36, 0.18, size=9, color=SUBTLE, align=PP_ALIGN.CENTER)


def build_deck() -> Presentation:
    root = Path(__file__).resolve().parents[1]
    base = root / "yelp_text_to_sql"
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slides: list[dict[str, object]] = [
        {
            "kicker": "SilkByte X Query",
            "title": "Big Data Analytics Project Summary",
            "subtitle": "Compact presentation covering the required project sections without overloading the deck.",
            "chips": ["Yelp Dataset", "HDFS", "Hive", "PySpark", "Zeppelin", "Text-to-SQL"],
            "bullets": [
                "End-to-end system from raw Yelp JSON ingestion to conversational analytics.",
                "Designed for explainable querying: question -> SQL -> result table -> chart.",
                "Built as both a big data coursework pipeline and a product-style demo.",
            ],
            "image": "assets/data_journey/generated/closing_value_chain_v2.png",
            "caption": "Project flow from storage and analytics to conversational insight",
        },
        {
            "kicker": "General Project Description",
            "title": "What The Project Does",
            "subtitle": "A unified system that converts large Yelp data into accessible business insight.",
            "chips": ["Storage", "Warehousing", "Processing", "Visualization", "AI Querying"],
            "bullets": [
                "Ingested semi-structured Yelp JSON into a scalable distributed environment.",
                "Structured the data in Hive and processed analytical tasks with PySpark.",
                "Used Zeppelin for notebook-based visual reporting and validation.",
                "Added Query by SilkByteX as the conversational layer for natural-language querying.",
            ],
            "image": "assets/data_journey/generated/etl_pipeline_flow.png",
            "caption": "Raw data -> warehouse -> analytics -> product layer",
        },
        {
            "kicker": "Team Modules",
            "title": "Member Responsibilities",
            "subtitle": "Module ownership was split by analysis domain, then integrated into one final system.",
            "chips": ["Team Division", "Ownership", "Integration"],
            "bullets_left": [
                "Sabah Laajaje: Business analysis and data enrichment.",
                "Focused on categories, ratings, locations, top-performing businesses, and external-context hypotheses.",
                "Fahim: User analysis and enrichment.",
                "Focused on growth, engagement, top reviewers, elite behavior, and external influence patterns.",
            ],
            "bullets_right": [
                "Mobashir Sifat: Review analysis and integration support.",
                "Focused on review trends, sentiment-oriented interpretation, and presentation/product synthesis.",
                "Shared team work: ETL workflow, architecture integration, Zeppelin outputs, and final delivery alignment.",
            ],
        },
        {
            "kicker": "System Architecture",
            "title": "Three-Layer Architecture",
            "subtitle": "The product is organized as presentation, service, and data-execution layers.",
            "chips": ["Streamlit UI", "Pipeline", "LLM", "Hive/Spark", "Schema Injection"],
            "bullets": [
                "Presentation layer in ui.py captures questions and renders SQL, tables, and charts.",
                "Service layer in app.py, pipeline.py, and sql_generation.py handles orchestration and self-correction.",
                "Data layer in database.py and schema_definitions.py connects the product to Hive/Spark execution.",
                "Schema-aware prompting improves SQL correctness before execution.",
            ],
            "image": "assets/data_journey/generated/query_architecture_blueprint.png",
            "caption": "Architecture used in the final Query by SilkByteX application",
        },
        {
            "kicker": "Project Highlights",
            "title": "What Makes The Project Strong",
            "subtitle": "The system combines engineering depth with a usable interface for decision-making.",
            "chips": ["End-to-End", "Scalable", "Explainable", "Interactive"],
            "bullets": [
                "Full pipeline from raw JSON to business-ready analytics.",
                "Multiple analysis domains: business, user, review, rating, and check-in.",
                "Conversational Text-to-SQL interface on top of a real backend workflow.",
                "Explainable outputs through SQL trace, tables, and chart views.",
                "Data enrichment direction extends the system beyond a single dataset.",
            ],
            "image": "assets/data_journey/generated/product_transition_storyboard.png",
            "caption": "From engineering pipeline to user-facing product experience",
        },
        {
            "kicker": "Data Insights",
            "title": "Key Insights 1 To 4",
            "subtitle": "Representative insight set drawn from the analysis scope in the project materials.",
            "chips": ["Business", "User", "Review", "Rating"],
            "bullets_left": [
                "Merchant concentration is uneven across cities and categories, revealing localized market dominance.",
                "Elite and high-activity users contribute a disproportionate share of review influence.",
                "Text patterns in reviews help explain why star ratings rise or fall, not just what the final score is.",
                "High review volume does not always mean consistently high satisfaction.",
            ],
            "bullets_right": [
                "Business category mix affects both visibility and rating behavior.",
                "User participation quality matters more than raw activity count alone.",
                "Pain-point language in low-star reviews exposes operational issues hidden by averages.",
                "Comparing categories and cities reveals more than looking at global averages.",
            ],
            "image": "assets/data_journey/ppt/image20.png",
            "caption": "Category and quality patterns visible in the analysis outputs",
        },
        {
            "kicker": "Data Insights",
            "title": "Key Insights 5 To 8",
            "subtitle": "Additional insights connect traffic behavior, timing, and enrichment opportunities.",
            "chips": ["Check-in", "Timing", "Enrichment", "Decision Support"],
            "bullets_left": [
                "Check-in activity complements reviews by showing physical demand, not only digital feedback.",
                "Weekday versus weekend patterns can change rating and traffic behavior.",
                "Temporal views help distinguish stable businesses from short-term spikes.",
                "Combining ratings, reviews, and check-ins creates stronger ranking logic than any single metric.",
            ],
            "bullets_right": [
                "External context such as weather or location conditions can refine internal interpretations.",
                "Cross-validation reduces the risk of drawing conclusions from one source alone.",
                "Conversational querying lowers the barrier for non-technical users to explore these insights.",
                "The product layer turns analysis outputs into reusable decision support.",
            ],
            "image": "assets/data_journey/ppt/image28.png",
            "caption": "Traffic behavior and comprehensive analysis reinforce final insights",
        },
        {
            "kicker": "Lessons Learned",
            "title": "What The Team Learned",
            "subtitle": "The project improved both technical depth and collaboration discipline.",
            "chips": ["Distributed Systems", "Analytics", "Collaboration", "Product Thinking"],
            "bullets": [
                "Clean schema design and data preparation are essential before advanced analytics.",
                "Distributed tools like HDFS, Hive, Spark, and Zeppelin work best as one coordinated workflow.",
                "Reliable analytics products need both backend rigor and frontend clarity.",
                "Team task division helps velocity, but integration quality determines the final result.",
                "AI querying is most valuable when it remains grounded in real schema and real outputs.",
            ],
            "image": "assets/data_journey/generated/closing_value_chain_v2.png",
            "caption": "Technical learning and product thinking connected across the full journey",
        },
        {
            "kicker": "Git Collaboration",
            "title": "Repository And Commit History",
            "subtitle": "Version control was part of the required workflow, but the current workspace snapshot does not include local .git metadata.",
            "chips": ["GitHub Workflow", "Collaboration", "Traceability"],
            "bullets": [
                "The provided team deck explicitly documents a GitHub repository overview and commit tracking as part of collaboration.",
                "Repository organization separated analysis modules and supported clearer task ownership.",
                "Git history was used to coordinate revisions and track ongoing progress across modules.",
                "Exact commit counts and per-member commit statistics are not recoverable from this exported workspace because no .git directory is present.",
            ],
            "image": "assets/data_journey/generated/agenda_timeline.png",
            "caption": "Collaboration rhythm and versioned project progress",
        },
        {
            "kicker": "Closing",
            "title": "From Raw Data To Conversational Intelligence",
            "subtitle": "A compact summary of the project’s engineering value and presentation-ready outcome.",
            "chips": ["ETL", "Analytics", "Insights", "AI Product"],
            "bullets": [
                "Built a real big data workflow, not isolated scripts.",
                "Generated actionable insight across multiple analytical domains.",
                "Translated the backend pipeline into a user-facing conversational product.",
                "Delivered a presentation-ready system with clear ownership, architecture, and lessons learned.",
            ],
            "image": "assets/data_journey/generated/demo_queries_showcase.png",
            "caption": "The final system turns questions into explainable outputs",
        },
    ]

    for index, data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        decorate(slide, base)
        add_card(slide, 0.48, 0.38, 12.35, 6.72)
        chip_bottom = title_block(
            slide,
            str(data["kicker"]),
            str(data["title"]),
            str(data["subtitle"]),
            list(data.get("chips", [])),
        )

        if "bullets_left" in data and "bullets_right" in data:
            add_card(slide, 0.78, chip_bottom + 0.1, 5.85, 3.95, fill=CARD_ALT)
            add_card(slide, 6.72, chip_bottom + 0.1, 5.35, 3.95, fill=CARD_ALT)
            add_bullet_lines(slide, list(data["bullets_left"]), 1.0, chip_bottom + 0.35, 5.4, 3.4, size=13)
            add_bullet_lines(slide, list(data["bullets_right"]), 6.95, chip_bottom + 0.35, 4.9, 3.25, size=13)
        elif "bullets_left" in data:
            add_card(slide, 0.78, chip_bottom + 0.1, 5.75, 3.95, fill=CARD_ALT)
            add_card(slide, 6.72, chip_bottom + 0.1, 5.55, 3.95, fill=CARD_ALT)
            add_bullet_lines(slide, list(data["bullets_left"]), 1.0, chip_bottom + 0.35, 5.3, 3.35, size=13)
            add_bullet_lines(slide, list(data["bullets_right"]), 6.95, chip_bottom + 0.35, 5.1, 3.35, size=13)
        else:
            add_card(slide, 0.78, chip_bottom + 0.08, 5.55, 3.98, fill=CARD_ALT)
            add_bullet_lines(slide, list(data["bullets"]), 1.0, chip_bottom + 0.32, 5.1, 3.45, size=13)
            add_visual_card(
                slide,
                base,
                data.get("image"),
                str(data.get("caption", "")),
                6.56,
                chip_bottom + 0.08,
                5.58,
                3.98,
            )

        add_textbox(slide, f"{index + 1:02d}", 12.1, 7.02, 0.4, 0.2, size=9, color=SUBTLE, align=PP_ALIGN.RIGHT)

    return prs


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out_path = root / "output" / "SILKBYTE_X_QUERY_Required_Compact.pptx"
    prs = build_deck()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Created: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
